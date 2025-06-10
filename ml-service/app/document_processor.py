import hashlib
from typing import Dict, List, Tuple
import logging
from pathlib import Path

# Document processing imports
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl
from bs4 import BeautifulSoup
import chromadb
from chromadb.config import Settings

logger = logging.getLogger(__name__)

class DocumentProcessor:
    def __init__(self, chroma_persist_directory: str = "./chroma_db"):
        """Initialize document processor with Chroma vector database"""
        self.chroma_persist_directory = chroma_persist_directory
        
        # Initialize Chroma client
        self.chroma_client = chromadb.PersistentClient(
            path=chroma_persist_directory,
            settings=Settings(
                anonymized_telemetry=False,
                allow_reset=True
            )
        )
        
        # Create or get collection
        self.collection_name = "documents"
        try:
            self.collection = self.chroma_client.get_collection(name=self.collection_name)
        except:
            self.collection = self.chroma_client.create_collection(
                name=self.collection_name,
                metadata={"hnsw:space": "cosine"}
            )
        
        logger.info(f"Chroma collection '{self.collection_name}' initialized")

    def extract_text_from_pdf(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from PDF file"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                metadata = {
                    "pages": len(pdf_reader.pages),
                    "title": pdf_reader.metadata.get('/Title', '') if pdf_reader.metadata else '',
                    "author": pdf_reader.metadata.get('/Author', '') if pdf_reader.metadata else '',
                    "subject": pdf_reader.metadata.get('/Subject', '') if pdf_reader.metadata else ''
                }
                
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    text += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
                
                return text.strip(), metadata
        except Exception as e:
            logger.error(f"Error extracting PDF text: {str(e)}")
            raise

    def extract_text_from_docx(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from Word document"""
        try:
            doc = Document(file_path)
            text = ""
            metadata = {
                "paragraphs": len(doc.paragraphs),
                "title": doc.core_properties.title or '',
                "author": doc.core_properties.author or '',
                "subject": doc.core_properties.subject or ''
            }
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            return text.strip(), metadata
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {str(e)}")
            raise

    def extract_text_from_pptx(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            text = ""
            slide_count = len(prs.slides)
            
            metadata = {
                "slides": slide_count,
                "title": prs.core_properties.title or '',
                "author": prs.core_properties.author or '',
                "subject": prs.core_properties.subject or ''
            }
            
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n--- Slide {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
            
            return text.strip(), metadata
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {str(e)}")
            raise

    def extract_text_from_xlsx(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from Excel file"""
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text = ""
            metadata = {
                "worksheets": len(workbook.worksheets),
                "worksheet_names": [ws.title for ws in workbook.worksheets]
            }
            
            for worksheet in workbook.worksheets:
                text += f"\n--- Worksheet: {worksheet.title} ---\n"
                for row in worksheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
            
            return text.strip(), metadata
        except Exception as e:
            logger.error(f"Error extracting XLSX text: {str(e)}")
            raise

    def extract_text_from_txt(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
                metadata = {
                    "encoding": "utf-8",
                    "lines": len(text.splitlines())
                }
                return text, metadata
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1') as file:
                text = file.read()
                metadata = {
                    "encoding": "latin-1",
                    "lines": len(text.splitlines())
                }
                return text, metadata
        except Exception as e:
            logger.error(f"Error extracting TXT text: {str(e)}")
            raise

    def process_document(self, file_path: str, file_name: str, file_id: int) -> Dict:
        """Process document and extract text based on file type"""
        file_extension = Path(file_path).suffix.lower()
        
        extractors = {
            '.pdf': self.extract_text_from_pdf,
            '.docx': self.extract_text_from_docx,
            '.doc': self.extract_text_from_docx,  # Note: limited support for .doc
            '.pptx': self.extract_text_from_pptx,
            '.xlsx': self.extract_text_from_xlsx,
            '.xls': self.extract_text_from_xlsx,  # Note: limited support for .xls
            '.txt': self.extract_text_from_txt,
            '.md': self.extract_text_from_txt,
            '.rtf': self.extract_text_from_txt
        }
        
        if file_extension not in extractors:
            raise ValueError(f"Unsupported file type: {file_extension}")
        
        try:
            extracted_text, doc_metadata = extractors[file_extension](file_path)
            
            # Create document hash for deduplication
            doc_hash = hashlib.md5(extracted_text.encode()).hexdigest()
            
            return {
                "file_id": file_id,
                "file_name": file_name,
                "file_path": file_path,
                "extracted_text": extracted_text,
                "doc_hash": doc_hash,
                "metadata": doc_metadata,
                "file_type": file_extension,
                "text_length": len(extracted_text)
            }
        except Exception as e:
            logger.error(f"Error processing document {file_name}: {str(e)}")
            raise

    def chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks for better embedding"""
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Try to break at sentence or paragraph boundaries
            if end < len(text):
                # Look for sentence endings
                for i in range(end, start + chunk_size - 100, -1):
                    if text[i] in '.!?\n':
                        end = i + 1
                        break
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end - overlap
            if start >= len(text):
                break
        
        return chunks

    async def store_in_vector_db(self, processed_doc: Dict, embeddings: List[List[float]]) -> bool:
        """Store document chunks and embeddings in Chroma vector database"""
        try:
            # Chunk the text
            chunks = self.chunk_text(processed_doc["extracted_text"])
            
            # Prepare data for Chroma
            ids = []
            documents = []
            metadatas = []
            chunk_embeddings = []
            
            for i, chunk in enumerate(chunks):
                chunk_id = f"{processed_doc['file_id']}_{i}"
                ids.append(chunk_id)
                documents.append(chunk)
                
                metadata = {
                    "file_id": processed_doc["file_id"],
                    "file_name": processed_doc["file_name"],
                    "file_path": processed_doc["file_path"],
                    "chunk_index": i,
                    "total_chunks": len(chunks),
                    "file_type": processed_doc["file_type"],
                    "doc_hash": processed_doc["doc_hash"],
                    **processed_doc["metadata"]
                }
                metadatas.append(metadata)
                
                # Use corresponding embedding or first one if single embedding
                if i < len(embeddings):
                    chunk_embeddings.append(embeddings[i])
                else:
                    chunk_embeddings.append(embeddings[0])
            
            # Add to Chroma collection
            self.collection.add(
                ids=ids,
                documents=documents,
                metadatas=metadatas,
                embeddings=chunk_embeddings
            )
            
            logger.info(f"Stored {len(chunks)} chunks for document {processed_doc['file_name']} in vector DB")
            return True
            
        except Exception as e:
            logger.error(f"Error storing in vector DB: {str(e)}")
            return False

    def search_similar_documents(self, query_embedding: List[float], n_results: int = 5) -> Dict:
        """Search for similar documents using vector similarity"""
        try:
            results = self.collection.query(
                query_embeddings=[query_embedding],
                n_results=n_results
            )
            
            return {
                "ids": results["ids"][0],
                "documents": results["documents"][0],
                "metadatas": results["metadatas"][0],
                "distances": results["distances"][0] if "distances" in results else []
            }
        except Exception as e:
            logger.error(f"Error searching vector DB: {str(e)}")
            return {"ids": [], "documents": [], "metadatas": [], "distances": []}

    def delete_document_from_vector_db(self, file_id: int) -> bool:
        """Delete all chunks of a document from vector database"""
        try:
            # Get all chunks for this file_id
            results = self.collection.get(
                where={"file_id": file_id}
            )
            
            if results["ids"]:
                self.collection.delete(ids=results["ids"])
                logger.info(f"Deleted {len(results['ids'])} chunks for file_id {file_id} from vector DB")
            
            return True
        except Exception as e:
            logger.error(f"Error deleting from vector DB: {str(e)}")
            return False